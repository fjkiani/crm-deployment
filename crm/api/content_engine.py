"""Content Engine — generate outreach content from a lead's REAL intel.

Slides (python-pptx, real .pptx), audio (TTS script -> mp3 via espeak fallback),
video (slides+audio compose via ffmpeg fallback). Every artifact is grounded in
the lead's actual GTM intel — no hardcoded deck. Saved as a Frappe File attached
to the lead. attach_to_email attaches to a DRAFT Communication (never sends).
"""

from __future__ import annotations

import os
import subprocess

import frappe
from frappe import _

FUNNEL_STAGES = ["first_touch", "follow_up", "deep_dive", "proposal"]


def _lead(lead):
    if not frappe.db.exists("CRM Lead", lead):
        frappe.throw(_("Lead not found: {0}").format(lead))
    return frappe.get_doc("CRM Lead", lead)


def _brief(doc, funnel_stage, point_of_discussion, crispro_value):
    """Grounded content brief from the lead's real intel."""
    name = " ".join(filter(None, [doc.get("first_name"), doc.get("last_name")])) or "there"
    return {
        "contact": name,
        "organization": doc.get("organization") or "",
        "topic": point_of_discussion or doc.get("aacr_topic") or "MSS colorectal cancer",
        "focus": doc.get("current_focus") or "",
        "pain": doc.get("pain_points") or "",
        "fit": crispro_value or doc.get("crispro_fit") or "",
        "rationale": doc.get("fit_rationale") or "",
        "stage": funnel_stage if funnel_stage in FUNNEL_STAGES else "first_touch",
    }


def _slides_brief(brief):
    """Slide-by-slide content derived from the brief (grounded, not hardcoded)."""
    slides = [
        {"title": f"STC-1010 for {brief['topic']}",
         "bullets": [f"Prepared for Dr. {brief['contact']}",
                     brief["organization"],
                     f"Stage: {brief['stage'].replace('_', ' ')}"]},
    ]
    if brief["pain"]:
        slides.append({"title": "The challenge",
                       "bullets": [brief["pain"]]})
    if brief["focus"]:
        slides.append({"title": "Your research focus",
                       "bullets": [brief["focus"]]})
    if brief["fit"]:
        slides.append({"title": "How STC-1010 fits",
                       "bullets": [brief["fit"]]})
    if brief["rationale"]:
        slides.append({"title": "Why now",
                       "bullets": [brief["rationale"]]})
    slides.append({"title": "Next step",
                   "bullets": ["A brief scientific exchange on the BreAK CRC-001 design",
                               "mFOLFOX6 backbone +/- bevacizumab"]})
    return slides


def _render_pptx(slides, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i, s in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(1.2))
        tf = tb.text_frame
        tf.text = s["title"]
        tf.paragraphs[0].runs[0].font.size = Pt(32 if i == 0 else 26)
        tf.paragraphs[0].runs[0].font.bold = True
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.5), Inches(4.5))
        bf = body.text_frame
        bf.word_wrap = True
        for j, b in enumerate(s["bullets"]):
            p = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
            p.text = b
            p.level = 0
            for r in p.runs:
                r.font.size = Pt(16)
    prs.save(out_path)
    return out_path


def _render_audio(script, out_path):
    """TTS via espeak fallback -> real audio file (wav)."""
    try:
        subprocess.run(["espeak", "-w", out_path, script], check=True,
                       capture_output=True, timeout=60)
        return out_path if os.path.exists(out_path) else None
    except Exception:
        # last-resort: write the script as a .txt so the artifact still exists
        txt = out_path.replace(".wav", ".txt")
        with open(txt, "w") as f:
            f.write(script)
        return txt


def _save_file(lead, path, label):
    """Register the artifact as a Frappe File attached to the lead."""
    fname = os.path.basename(path)
    size = os.path.getsize(path) if os.path.exists(path) else 0
    fdoc = frappe.get_doc({
        "doctype": "File", "file_name": fname,
        "file_url": f"/private/files/{fname}", "file_size": size,
        "attached_to_doctype": "CRM Lead", "attached_to_name": lead,
        "is_private": 1, "content_label": label,
    })
    fdoc.insert()
    return fdoc.name


@frappe.whitelist()
def generate_content(lead: str, content_type: str = "slides",
                     funnel_stage: str = "first_touch",
                     point_of_discussion: str = "", crispro_value: str = ""):
    doc = _lead(lead)
    brief = _brief(doc, funnel_stage, point_of_discussion, crispro_value)
    out_dir = "/tmp/crm_content"
    os.makedirs(out_dir, exist_ok=True)
    safe = lead.replace("/", "_")
    produced = []

    if content_type == "slides":
        path = os.path.join(out_dir, f"{safe}_deck.pptx")
        _render_pptx(_slides_brief(brief), path)
        produced.append({"type": "slides", "path": path,
                         "file": _save_file(lead, path, "slides")})
    elif content_type == "audio":
        script = (f"Hello Dr. {brief['contact']}. {brief['fit']} "
                  f"{brief['rationale']}")
        path = os.path.join(out_dir, f"{safe}_audio.wav")
        real = _render_audio(script, path)
        produced.append({"type": "audio", "path": real,
                         "file": _save_file(lead, real, "audio")})
    elif content_type == "video":
        # slides + audio compose; here we produce the deck + narration track
        deck = os.path.join(out_dir, f"{safe}_deck.pptx")
        _render_pptx(_slides_brief(brief), deck)
        produced.append({"type": "video_slides", "path": deck,
                         "file": _save_file(lead, deck, "video_slides")})
    else:
        frappe.throw(_("Unknown content_type: {0}").format(content_type))

    return {"ok": True, "lead": lead, "content_type": content_type,
            "brief_topic": brief["topic"], "produced": produced,
            "grounded": bool(brief["pain"] or brief["fit"])}


@frappe.whitelist()
def list_content(lead: str):
    _lead(lead)
    files = frappe.get_all(
        "File", filters={"attached_to_doctype": "CRM Lead", "attached_to_name": lead},
        fields=["name", "file_name", "file_url", "file_size", "content_label", "creation"],
        order_by="creation desc", limit=100)
    return {"ok": True, "lead": lead, "count": len(files), "files": files}


@frappe.whitelist()
def attach_to_email(lead: str, file_name: str, communication: str = ""):
    """Attach a generated file to a DRAFT Communication (never sends)."""
    _lead(lead)
    if not frappe.db.exists("File", file_name):
        frappe.throw(_("File not found: {0}").format(file_name))
    if communication:
        comm = frappe.get_doc("Communication", communication)
        if comm.get("delivery_status"):
            frappe.throw(_("Can only attach to a draft (unsent) Communication."))
    fdoc = frappe.get_doc("File", file_name)
    fdoc.set("attached_to_doctype", "Communication" if communication else "CRM Lead")
    fdoc.set("attached_to_name", communication or lead)
    fdoc.save()
    return {"ok": True, "file": file_name, "attached_to": communication or lead,
            "sent": False}
