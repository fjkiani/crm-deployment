"""Real NotebookLM / Gemini Notebook content engine.

Provider-abstracted: ONE interface, THREE real backends. There is no local
synthesis fallback. When a backend's credential is absent, the provider raises
`NotebookLMCredentialError` naming exactly what is required and how to supply it.
It NEVER silently falls back to espeak / python-pptx. That masquerade is gone.

Backends
--------
gemini      Official Gemini API (generativelanguage). Grounded narrative via a
            Gemini text model + a real two-host audio overview via Gemini 2.5
            multi-speaker TTS. Credential: GEMINI_API_KEY (or GOOGLE_API_KEY).
            kinds: audio (real .wav), slides (Gemini-authored deck, rendered
            locally and labelled gemini-native — NOT literally NotebookLM).

enterprise  Official NotebookLM Enterprise REST (Discovery Engine / Agentspace):
            notebooks.create -> sources.batchCreate -> notebooks.audioOverviews
            .create. Credential: OAuth bearer token + GCP project number +
            location + a Gemini Enterprise license. kinds: audio ONLY — the
            official API exposes no slide-deck or video endpoint.

unofficial  The actual NotebookLM / Gemini Notebook product driven through the
            real `notebooklm` CLI (notebooklm-py). Produces LITERAL NotebookLM
            slide decks (.pptx/.pdf), video (.mp4) and audio (.mp3).
            Credential: a logged-in Google session (storage_state.json from
            `notebooklm login`). Unofficial / ToS-gray; may break on Google
            service changes. kinds: slides, audio, video.

Every function here is credential-gated and fail-loud. Nothing is stubbed,
cached, or synthesised locally as a stand-in for a real provider result.
"""
from __future__ import annotations

import os
import json
import shutil
import struct
import subprocess
import wave

# ---- model / endpoint constants (real) --------------------------------------
GEMINI_TEXT_MODEL = "gemini-2.5-flash"
GEMINI_TTS_MODEL = "gemini-2.5-pro-tts"          # multi-speaker studio TTS
GEMINI_TTS_SAMPLE_RATE = 24000                    # Gemini TTS PCM output rate
ENTERPRISE_API_VERSION = "v1alpha"
ENTERPRISE_HOST_TMPL = "https://{loc}-discoveryengine.googleapis.com"
ENTERPRISE_MULTIREGIONS = {"us", "eu", "global"}

# kinds a given provider can actually deliver
PROVIDER_KINDS = {
    "gemini": {"audio", "slides"},
    "enterprise": {"audio"},
    "unofficial": {"slides", "audio", "video"},
}


# ---- errors -----------------------------------------------------------------
class NotebookLMError(Exception):
    """Base error for the NotebookLM engine."""


class NotebookLMCredentialError(NotebookLMError):
    """A provider was asked to run without the credential it requires."""

    def __init__(self, provider, missing, how):
        self.provider = provider
        self.missing = list(missing)
        self.how = how
        super().__init__(
            "[{p}] missing credential(s): {m}. {how}".format(
                p=provider, m=", ".join(self.missing), how=how
            )
        )


class NotebookLMUnsupportedKind(NotebookLMError):
    """A provider cannot produce the requested content kind."""


# ---- credential discovery (presence only) ----------------------------------
def _gemini_key():
    return os.environ.get("GEMINI_API_KEY") or os.environ.get("GOOGLE_API_KEY")


def _enterprise_creds():
    return {
        "token": os.environ.get("NOTEBOOKLM_OAUTH_TOKEN"),
        "project": os.environ.get("NOTEBOOKLM_GCP_PROJECT"),
        "location": os.environ.get("NOTEBOOKLM_LOCATION", "us"),
    }


def _unofficial_cli():
    return os.environ.get("NOTEBOOKLM_CLI") or shutil.which("notebooklm")


def _unofficial_storage():
    override = os.environ.get("NOTEBOOKLM_STORAGE")
    if override and os.path.exists(override):
        return override
    default = os.path.expanduser("~/.notebooklm/profiles/default/storage_state.json")
    return default if os.path.exists(default) else None


def available_providers():
    """Report which providers are live-capable in the current environment.

    Returns a dict {provider: {"live": bool, "missing": [...], "kinds": [...]}}.
    Presence checks only — never returns secret values.
    """
    out = {}

    out["gemini"] = {
        "live": bool(_gemini_key()),
        "missing": [] if _gemini_key() else ["GEMINI_API_KEY"],
        "kinds": sorted(PROVIDER_KINDS["gemini"]),
    }

    ent = _enterprise_creds()
    ent_missing = [k for k in ("token", "project") if not ent.get(k)]
    out["enterprise"] = {
        "live": not ent_missing,
        "missing": ["NOTEBOOKLM_OAUTH_TOKEN" if m == "token" else "NOTEBOOKLM_GCP_PROJECT"
                    for m in ent_missing],
        "kinds": sorted(PROVIDER_KINDS["enterprise"]),
    }

    cli, storage = _unofficial_cli(), _unofficial_storage()
    umiss = []
    if not cli:
        umiss.append("notebooklm CLI (pip install notebooklm-py)")
    if not storage:
        umiss.append("Google session (run: notebooklm login)")
    out["unofficial"] = {
        "live": not umiss,
        "missing": umiss,
        "kinds": sorted(PROVIDER_KINDS["unofficial"]),
    }
    return out


# ---- grounding --------------------------------------------------------------
FUNNEL_STAGES = ("first_touch", "follow_up", "deep_dive", "proposal")


def build_sources(brief):
    """Turn grounded lead intel into a single source document.

    `brief` keys (all optional except contact/company): contact, company, title,
    topic, current_focus, pain_points, crispro_fit, fit_rationale, funnel_stage,
    point_of_discussion, crispro_value.
    This text is what gets fed to NotebookLM as the *source* (or to Gemini as the
    grounding context). Nothing is invented here — only the lead's own intel.
    """
    b = brief or {}
    lines = []
    who = b.get("contact") or "the prospect"
    company = b.get("company") or ""
    lines.append("# Engagement source dossier: {who}{comp}".format(
        who=who, comp=(" — " + company) if company else ""))
    if b.get("title"):
        lines.append("Role: {}".format(b["title"]))
    if b.get("topic"):
        lines.append("\n## Their stated topic\n{}".format(b["topic"]))
    if b.get("current_focus"):
        lines.append("\n## Their current focus\n{}".format(b["current_focus"]))
    if b.get("pain_points"):
        lines.append("\n## Their pain points\n{}".format(b["pain_points"]))
    if b.get("crispro_fit"):
        lines.append("\n## CrisPRO fit\n{}".format(b["crispro_fit"]))
    if b.get("fit_rationale"):
        lines.append("\n## Fit rationale\n{}".format(b["fit_rationale"]))
    if b.get("point_of_discussion"):
        lines.append("\n## Point of discussion for this piece\n{}".format(b["point_of_discussion"]))
    if b.get("crispro_value"):
        lines.append("\n## CrisPRO value to emphasise\n{}".format(b["crispro_value"]))
    # Require at least one substantive intel field — never ground on boilerplate.
    substantive = ("topic", "current_focus", "pain_points", "crispro_fit", "fit_rationale")
    if not any((b.get(k) or "").strip() for k in substantive):
        raise NotebookLMError(
            "Refusing to generate: lead intel is too thin to ground content "
            "(need topic / focus / pain points / fit). Enrich the lead first."
        )
    stage = b.get("funnel_stage") or "first_touch"
    lines.append("\n## Funnel stage\n{}".format(stage))
    return "\n".join(lines).strip()


def _instructions(kind, brief):
    b = brief or {}
    stage = b.get("focus") or b.get("point_of_discussion") or b.get("funnel_stage") or "first_touch"
    who = b.get("contact") or "the prospect"
    return (
        "Create a {stage} {kind} for {who}. Ground every claim strictly in the "
        "provided source dossier — their topic, focus, pain points and the "
        "CrisPRO fit. Do not invent data. Lead with their problem, then how "
        "STC-1010 addresses it.".format(stage=stage, kind=kind, who=who)
    )


# ---- pure request builders (network-free, unit-testable) --------------------
def enterprise_notebook_request(project, location, title, token):
    if location not in ENTERPRISE_MULTIREGIONS:
        raise NotebookLMError("location must be one of {}".format(ENTERPRISE_MULTIREGIONS))
    host = ENTERPRISE_HOST_TMPL.format(loc=location)
    url = "{host}/{v}/projects/{proj}/locations/{loc}/notebooks".format(
        host=host, v=ENTERPRISE_API_VERSION, proj=project, loc=location)
    headers = {"Authorization": "Bearer {}".format(token), "Content-Type": "application/json"}
    return ("POST", url, headers, {"title": title})


def enterprise_source_request(project, location, notebook_id, sources, token):
    host = ENTERPRISE_HOST_TMPL.format(loc=location)
    url = "{host}/{v}/projects/{proj}/locations/{loc}/notebooks/{nb}/sources:batchCreate".format(
        host=host, v=ENTERPRISE_API_VERSION, proj=project, loc=location, nb=notebook_id)
    headers = {"Authorization": "Bearer {}".format(token), "Content-Type": "application/json"}
    return ("POST", url, headers, {"userContents": sources})


def enterprise_audio_request(project, location, notebook_id, source_ids, focus, language, token):
    host = ENTERPRISE_HOST_TMPL.format(loc=location)
    url = "{host}/{v}/projects/{proj}/locations/{loc}/notebooks/{nb}/audioOverviews".format(
        host=host, v=ENTERPRISE_API_VERSION, proj=project, loc=location, nb=notebook_id)
    headers = {"Authorization": "Bearer {}".format(token), "Content-Type": "application/json"}
    body = {"episodeFocus": focus, "languageCode": language}
    if source_ids:
        body["sourceIds"] = [{"id": s} for s in source_ids]
    return ("POST", url, headers, body)


def unofficial_argv(cli, storage, subcmd, notebook_id=None, description=None,
                    fmt=None, length=None, source_text=None, source_title=None, title=None):
    """Build the exact `notebooklm` CLI argv for a step. No execution here."""
    base = [cli, "--storage", storage, "--quiet"]
    if subcmd == "create":
        return base + ["create", title or "CRM engagement", "--json"]
    if subcmd == "source-add":
        return base + ["source", "add", "-n", notebook_id, "--type", "text",
                       "--title", source_title or "Lead dossier", source_text or "", "--json"]
    if subcmd == "slide-deck":
        args = base + ["generate", "slide-deck", description or "", "-n", notebook_id,
                       "--format", fmt or "presenter", "--length", length or "default",
                       "--wait", "--json"]
        return args
    if subcmd == "audio":
        return base + ["generate", "audio", description or "", "-n", notebook_id,
                       "--format", fmt or "deep-dive", "--length", length or "default",
                       "--wait", "--json"]
    if subcmd == "video":
        return base + ["generate", "video", description or "", "-n", notebook_id,
                       "--format", fmt or "explainer", "--wait", "--json"]
    raise NotebookLMError("unknown unofficial subcmd: {}".format(subcmd))


def gemini_tts_config(speakers):
    """Build the REAL google-genai multi-speaker SpeechConfig. Imports the SDK."""
    from google.genai import types  # lazy: only needed when actually running
    speaker_cfgs = [
        types.SpeakerVoiceConfig(
            speaker=name,
            voice_config=types.VoiceConfig(
                prebuilt_voice_config=types.PrebuiltVoiceConfig(voice_name=voice)
            ),
        )
        for name, voice in speakers
    ]
    return types.GenerateContentConfig(
        response_modalities=["AUDIO"],
        speech_config=types.SpeechConfig(
            multi_speaker_voice_config=types.MultiSpeakerVoiceConfig(
                speaker_voice_configs=speaker_cfgs
            )
        ),
    )


# ---- audio helpers ----------------------------------------------------------
def _pcm_to_wav(pcm_bytes, path, rate=GEMINI_TTS_SAMPLE_RATE):
    with wave.open(path, "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)  # 16-bit
        w.setframerate(rate)
        w.writeframes(pcm_bytes)
    return path


# ---- providers --------------------------------------------------------------
class _Provider:
    name = "base"
    kinds = frozenset()

    def check(self):
        raise NotImplementedError

    def _assert_kind(self, kind):
        if kind not in self.kinds:
            raise NotebookLMUnsupportedKind(
                "[{p}] cannot produce '{k}'. Supported: {ks}.".format(
                    p=self.name, k=kind, ks=sorted(self.kinds))
            )

    def generate(self, kind, brief, out_dir, **opts):
        raise NotImplementedError


class GeminiProvider(_Provider):
    name = "gemini"
    kinds = frozenset({"audio", "slides"})

    def check(self):
        if not _gemini_key():
            raise NotebookLMCredentialError(
                "gemini", ["GEMINI_API_KEY"],
                "Set GEMINI_API_KEY (or GOOGLE_API_KEY) to a Google AI Studio "
                "key. Get one at https://aistudio.google.com/apikey.",
            )

    def _client(self):
        try:
            from google import genai  # lazy
        except Exception as exc:  # pragma: no cover - env dependent
            raise NotebookLMError(
                "google-genai not installed (pip install google-genai): {}".format(exc)
            )
        return genai.Client(api_key=_gemini_key())

    def generate(self, kind, brief, out_dir, **opts):
        self._assert_kind(kind)
        self.check()
        os.makedirs(out_dir, exist_ok=True)
        source = build_sources(brief)
        client = self._client()
        if kind == "audio":
            script = self._script(client, source, brief)
            audio = client.models.generate_content(
                model=GEMINI_TTS_MODEL,
                contents=script,
                config=gemini_tts_config([("Host", "Kore"), ("Guest", "Puck")]),
            )
            pcm = audio.candidates[0].content.parts[0].inline_data.data
            path = os.path.join(out_dir, "{}_overview.wav".format(_slug(brief)))
            _pcm_to_wav(pcm, path)
            return {"provider": self.name, "kind": kind, "path": path, "live": True,
                    "meta": {"model": GEMINI_TTS_MODEL, "speakers": 2, "grounded": True}}
        # slides: Gemini authors the deck content; rendered locally (labelled).
        deck = self._deck_content(client, source, brief)
        path = _render_pptx(deck, os.path.join(out_dir, "{}_deck.pptx".format(_slug(brief))))
        return {"provider": self.name, "kind": kind, "path": path, "live": True,
                "meta": {"model": GEMINI_TEXT_MODEL, "grounded": True,
                         "note": "Gemini-authored deck rendered locally (not literal NotebookLM)"}}

    def _script(self, client, source, brief):
        prompt = (
            "Write a two-host podcast transcript (~250 words) for a NotebookLM-style "
            "audio overview, grounded strictly in the source below. Two speakers "
            "named exactly 'Host' and 'Guest'. Prefix each line with 'Host:' or "
            "'Guest:'. Open with the prospect's problem, then how STC-1010 helps.\n\n"
            "SOURCE:\n" + source
        )
        resp = client.models.generate_content(model=GEMINI_TEXT_MODEL, contents=prompt)
        return resp.text

    def _deck_content(self, client, source, brief):
        from google.genai import types
        prompt = (
            "Author a 6-slide sales deck grounded strictly in the source. Return "
            "JSON: {\"slides\":[{\"title\":str,\"bullets\":[str,...]}]}. Slide 1 = "
            "their problem; last = clear next step. No invented data.\n\nSOURCE:\n" + source
        )
        resp = client.models.generate_content(
            model=GEMINI_TEXT_MODEL, contents=prompt,
            config=types.GenerateContentConfig(response_mime_type="application/json"),
        )
        data = json.loads(resp.text)
        return data.get("slides", [])


class EnterpriseProvider(_Provider):
    name = "enterprise"
    kinds = frozenset({"audio"})

    def check(self):
        c = _enterprise_creds()
        missing = []
        if not c.get("token"):
            missing.append("NOTEBOOKLM_OAUTH_TOKEN")
        if not c.get("project"):
            missing.append("NOTEBOOKLM_GCP_PROJECT")
        if missing:
            raise NotebookLMCredentialError(
                "enterprise", missing,
                "Provide a GCP project number and an OAuth token "
                "(gcloud auth print-access-token) with a Gemini Enterprise "
                "license + Discovery Engine API enabled.",
            )

    def generate(self, kind, brief, out_dir, **opts):
        self._assert_kind(kind)  # raises for slides/video with a clear message
        self.check()
        import requests  # real HTTP
        os.makedirs(out_dir, exist_ok=True)
        c = _enterprise_creds()
        proj, loc, token = c["project"], c["location"], c["token"]
        source = build_sources(brief)
        # 1. notebook
        m, url, h, body = enterprise_notebook_request(proj, loc, brief.get("contact") or "CRM", token)
        nb = requests.post(url, headers=h, json=body, timeout=60)
        nb.raise_for_status()
        nb_id = nb.json().get("notebookId") or nb.json().get("name", "").split("/")[-1]
        # 2. source
        m, url, h, body = enterprise_source_request(
            proj, loc, nb_id,
            [{"content": {"content": source, "mimeType": "text/plain"},
              "displayName": "Lead dossier"}], token)
        requests.post(url, headers=h, json=body, timeout=120).raise_for_status()
        # 3. audio overview
        focus = _instructions("audio overview", brief)
        m, url, h, body = enterprise_audio_request(proj, loc, nb_id, None, focus, "en", token)
        ao = requests.post(url, headers=h, json=body, timeout=120)
        ao.raise_for_status()
        return {"provider": self.name, "kind": kind, "notebook_id": nb_id, "live": True,
                "meta": {"operation": ao.json(), "grounded": True,
                         "note": "Poll the returned operation, then download the MP3."}}


class UnofficialProvider(_Provider):
    name = "unofficial"
    kinds = frozenset({"slides", "audio", "video"})

    def check(self):
        cli, storage = _unofficial_cli(), _unofficial_storage()
        missing = []
        if not cli:
            missing.append("notebooklm CLI")
        if not storage:
            missing.append("Google session (storage_state.json)")
        if missing:
            raise NotebookLMCredentialError(
                "unofficial", missing,
                "Install notebooklm-py and run `notebooklm login` (browser) to "
                "create a session; or point NOTEBOOKLM_STORAGE at a saved "
                "storage_state.json.",
            )

    def _run(self, argv):
        proc = subprocess.run(argv, capture_output=True, text=True)
        if proc.returncode != 0:
            raise NotebookLMError("notebooklm CLI failed ({}): {}".format(
                proc.returncode, (proc.stderr or proc.stdout)[:500]))
        try:
            return json.loads(proc.stdout or "{}")
        except json.JSONDecodeError:
            return {"raw": proc.stdout}

    def generate(self, kind, brief, out_dir, **opts):
        self._assert_kind(kind)
        self.check()
        os.makedirs(out_dir, exist_ok=True)
        cli, storage = _unofficial_cli(), _unofficial_storage()
        source = build_sources(brief)
        nb = self._run(unofficial_argv(cli, storage, "create",
                                       title=(brief.get("contact") or "CRM engagement")))
        nb_id = nb.get("id") or nb.get("notebookId") or nb.get("notebook", {}).get("id")
        if not nb_id:
            raise NotebookLMError("could not resolve notebook id from CLI output: {}".format(nb))
        self._run(unofficial_argv(cli, storage, "source-add", notebook_id=nb_id,
                                  source_text=source, source_title="Lead dossier"))
        subcmd = {"slides": "slide-deck", "audio": "audio", "video": "video"}[kind]
        res = self._run(unofficial_argv(cli, storage, subcmd, notebook_id=nb_id,
                                        description=_instructions(kind, brief)))
        path = res.get("path") or res.get("file") or (res.get("artifact") or {}).get("path")
        if not path or not os.path.exists(path):
            raise NotebookLMError(
                "notebooklm produced no downloadable file for '{}': {}".format(kind, res))
        return {"provider": self.name, "kind": kind, "path": path, "notebook_id": nb_id,
                "live": True, "meta": {"grounded": True, "note": "literal NotebookLM artifact"}}


_PROVIDERS = {"gemini": GeminiProvider, "enterprise": EnterpriseProvider, "unofficial": UnofficialProvider}

# preference order per kind (best real experience first)
_KIND_PREFERENCE = {
    "slides": ["unofficial", "gemini"],       # unofficial = literal NBLM deck
    "audio": ["gemini", "enterprise", "unofficial"],
    "video": ["unofficial"],
}


def _slug(brief):
    base = (brief or {}).get("slug") or (brief or {}).get("contact") or "engagement"
    return "".join(c if c.isalnum() else "_" for c in str(base))[:60]


def _render_pptx(slides, path):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for s in (slides or [{"title": "Overview", "bullets": []}]):
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(1)).text_frame
        tb.text = s.get("title", "")
        tb.paragraphs[0].font.size = Pt(30)
        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5)).text_frame
        for i, b in enumerate(s.get("bullets", []) or []):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = "• {}".format(b)
            p.font.size = Pt(18)
    prs.save(path)
    return path


def generate(kind, brief, provider="auto", out_dir="/tmp/crm_content", **opts):
    """Generate real content through a NotebookLM backend.

    kind: 'slides' | 'audio' | 'video'.
    provider: 'auto' (pick best live backend for the kind) or a specific name.
    Raises NotebookLMCredentialError (no live backend) or NotebookLMUnsupportedKind.
    Never returns synthesised/stub content.
    """
    if kind not in ("slides", "audio", "video"):
        raise NotebookLMError("kind must be slides|audio|video, got {}".format(kind))

    if provider != "auto":
        prov = _PROVIDERS.get(provider)
        if not prov:
            raise NotebookLMError("unknown provider {}".format(provider))
        return prov().generate(kind, brief, out_dir, **opts)

    avail = available_providers()
    order = _KIND_PREFERENCE.get(kind, list(_PROVIDERS))
    live = [p for p in order if avail.get(p, {}).get("live") and kind in PROVIDER_KINDS[p]]
    if live:
        return _PROVIDERS[live[0]]().generate(kind, brief, out_dir, **opts)

    # No live backend: fail loud with the exact unblock for each capable provider.
    capable = [p for p in order if kind in PROVIDER_KINDS[p]]
    hints = []
    for p in capable:
        try:
            _PROVIDERS[p]().check()
        except NotebookLMCredentialError as e:
            hints.append("{}: need {} — {}".format(p, ", ".join(e.missing), e.how))
    raise NotebookLMCredentialError(
        "auto", ["a Google credential for {}".format(kind)],
        "No live NotebookLM backend for '{}'. Options:\n  - ".format(kind)
        + "\n  - ".join(hints),
    )
